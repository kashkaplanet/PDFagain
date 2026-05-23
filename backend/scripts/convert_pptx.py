"""
PPTX to PDF converter using python-pptx + reportlab + Pillow.
No Microsoft Office installation required.
"""
import sys
import os
import io
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.colors import Color, black, white
from reportlab.lib.utils import ImageReader
from PIL import Image


def emu_to_points(emu):
    """Convert EMU (English Metric Units) to PDF points."""
    return emu / 914400 * 72


def pptx_color_to_reportlab(rgb_color, default=black):
    """Convert a python-pptx RGBColor to a reportlab Color."""
    if rgb_color is None:
        return default
    try:
        r = rgb_color[0] / 255.0
        g = rgb_color[1] / 255.0
        b = rgb_color[2] / 255.0
        return Color(r, g, b)
    except Exception:
        return default


def get_shape_fill_color(shape):
    """Try to extract fill color from a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fore_color = fill.fore_color
            if fore_color and fore_color.type is not None:
                return pptx_color_to_reportlab(fore_color.rgb)
    except Exception:
        pass
    return None


def get_slide_background_color(slide):
    """Try to get the slide background color."""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            fore_color = fill.fore_color
            if fore_color and fore_color.type is not None:
                return pptx_color_to_reportlab(fore_color.rgb)
    except Exception:
        pass
    return white


def draw_text_frame(c, text_frame, x, y, width, height, slide_height):
    """Draw text from a text frame onto the canvas."""
    # y in reportlab is from bottom, pptx is from top
    current_y = slide_height - y

    for paragraph in text_frame.paragraphs:
        # Get paragraph alignment
        alignment = paragraph.alignment

        # Build line text and track formatting
        line_parts = []
        for run in paragraph.runs:
            text = run.text
            if not text:
                continue

            font_size = 12  # default
            try:
                if run.font.size:
                    font_size = run.font.size.pt
            except Exception:
                pass

            font_color = black
            try:
                if run.font.color and run.font.color.rgb:
                    font_color = pptx_color_to_reportlab(run.font.color.rgb)
            except Exception:
                pass

            is_bold = False
            try:
                is_bold = run.font.bold
            except Exception:
                pass

            line_parts.append({
                'text': text,
                'size': font_size,
                'color': font_color,
                'bold': is_bold,
            })

        if not line_parts:
            # Empty paragraph - add spacing
            current_y -= 14
            continue

        # Determine max font size in this paragraph for line height
        max_size = max(p['size'] for p in line_parts)
        line_height = max_size * 1.3

        current_y -= line_height

        # Calculate text width for alignment
        full_text = ''.join(p['text'] for p in line_parts)

        # Draw each run
        text_x = x
        if alignment == PP_ALIGN.CENTER:
            # Approximate centering
            try:
                c.setFont("Helvetica", max_size)
                text_width = c.stringWidth(full_text, "Helvetica", max_size)
                text_x = x + (width - text_width) / 2
            except Exception:
                pass
        elif alignment == PP_ALIGN.RIGHT:
            try:
                c.setFont("Helvetica", max_size)
                text_width = c.stringWidth(full_text, "Helvetica", max_size)
                text_x = x + width - text_width
            except Exception:
                pass

        for part in line_parts:
            font_name = "Helvetica-Bold" if part['bold'] else "Helvetica"
            c.setFont(font_name, part['size'])
            c.setFillColor(part['color'])
            c.drawString(text_x, current_y, part['text'])
            text_x += c.stringWidth(part['text'], font_name, part['size'])


def convert_to_pdf(input_path, output_path):
    """Convert a PPTX file to PDF."""
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)

    if input_path.lower().endswith('.ppt'):
        try:
            from spire.presentation import Presentation as SpirePresentation, FileFormat
            p = SpirePresentation()
            p.LoadFromFile(input_path)
            p.SaveToFile(output_path, FileFormat.PDF)
            p.Dispose()
            print(f"Successfully converted {input_path} to {output_path} using Spire.Presentation")
            return
        except Exception as e:
            print(f"Failed to convert .ppt with Spire.Presentation: {e}", file=sys.stderr)
            sys.exit(1)

    prs = Presentation(input_path)

    # Get slide dimensions
    slide_width = emu_to_points(prs.slide_width)
    slide_height = emu_to_points(prs.slide_height)

    # Create PDF canvas with slide dimensions
    c = canvas.Canvas(output_path, pagesize=(slide_width, slide_height))

    for slide_num, slide in enumerate(prs.slides):
        # Draw background
        bg_color = get_slide_background_color(slide)
        c.setFillColor(bg_color)
        c.rect(0, 0, slide_width, slide_height, fill=1, stroke=0)

        # Process each shape
        for shape in slide.shapes:
            left = emu_to_points(shape.left) if shape.left else 0
            top = emu_to_points(shape.top) if shape.top else 0
            w = emu_to_points(shape.width) if shape.width else 0
            h = emu_to_points(shape.height) if shape.height else 0

            # Draw shape fill/background
            fill_color = get_shape_fill_color(shape)
            if fill_color:
                c.setFillColor(fill_color)
                rl_y = slide_height - top - h
                c.rect(left, rl_y, w, h, fill=1, stroke=0)

            # Handle images
            if shape.shape_type == 13:  # Picture
                try:
                    image_blob = shape.image.blob
                    image = Image.open(io.BytesIO(image_blob))
                    # Convert to RGB if necessary (handles RGBA, CMYK, etc.)
                    if image.mode in ('RGBA', 'LA', 'P'):
                        bg = Image.new('RGB', image.size, (255, 255, 255))
                        if image.mode == 'P':
                            image = image.convert('RGBA')
                        bg.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                        image = bg
                    elif image.mode != 'RGB':
                        image = image.convert('RGB')

                    img_reader = ImageReader(image)
                    rl_y = slide_height - top - h
                    c.drawImage(img_reader, left, rl_y, width=w, height=h,
                                preserveAspectRatio=True, mask='auto')
                except Exception as e:
                    print(f"Warning: Could not render image on slide {slide_num + 1}: {e}")

            # Handle tables
            elif shape.has_table:
                try:
                    table = shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)

                    if rows == 0 or cols == 0:
                        continue

                    cell_w = w / cols
                    cell_h = h / rows

                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_x = left + col_idx * cell_w
                            cell_y_top = top + row_idx * cell_h
                            rl_y = slide_height - cell_y_top - cell_h

                            # Draw cell border
                            c.setStrokeColor(black)
                            c.setLineWidth(0.5)
                            c.rect(cell_x, rl_y, cell_w, cell_h, fill=0, stroke=1)

                            # Draw cell text
                            if cell.text_frame:
                                c.setFillColor(black)
                                c.setFont("Helvetica", 9)
                                text = cell.text.strip()
                                if text:
                                    # Simple truncation for now
                                    max_chars = int(cell_w / 5)
                                    if len(text) > max_chars:
                                        text = text[:max_chars] + "..."
                                    c.drawString(cell_x + 3, rl_y + cell_h / 2 - 4, text)
                except Exception as e:
                    print(f"Warning: Could not render table on slide {slide_num + 1}: {e}")

            # Handle text frames
            elif shape.has_text_frame:
                try:
                    draw_text_frame(c, shape.text_frame, left, top, w, h, slide_height)
                except Exception as e:
                    print(f"Warning: Could not render text on slide {slide_num + 1}: {e}")

            # Handle grouped shapes
            elif hasattr(shape, 'shapes'):
                try:
                    for sub_shape in shape.shapes:
                        sub_left = emu_to_points(sub_shape.left) if sub_shape.left else 0
                        sub_top = emu_to_points(sub_shape.top) if sub_shape.top else 0
                        sub_w = emu_to_points(sub_shape.width) if sub_shape.width else 0
                        sub_h = emu_to_points(sub_shape.height) if sub_shape.height else 0

                        if sub_shape.has_text_frame:
                            draw_text_frame(c, sub_shape.text_frame,
                                            sub_left, sub_top, sub_w, sub_h, slide_height)
                except Exception as e:
                    print(f"Warning: Could not render group on slide {slide_num + 1}: {e}")

        # Add new page for next slide
        c.showPage()

    c.save()
    print(f"Successfully converted {input_path} to {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python convert_pptx.py <input_pptx> <output_pdf>")
        sys.exit(1)

    convert_to_pdf(sys.argv[1], sys.argv[2])
